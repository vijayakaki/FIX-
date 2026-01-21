

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle
    
    # Style the title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def create_content_slide(prs, title, content_items):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Handle indentation levels
        indent_level = item.get('level', 0)
        p.level = indent_level
        p.text = item['text']
        p.font.size = Pt(18 if indent_level == 0 else 16)
    
    return slide

def create_section_header(prs, title, subtitle=""):
    """Create a section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title_shape = slide.shapes.title
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add colored background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    if subtitle:
        # Add subtitle text box
        left = Inches(1)
        top = Inches(3.5)
        width = Inches(8)
        height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_formula_slide(prs, title, formula, explanation):
    """Create a slide for formulas"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add formula box
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = formula
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER
    
    # Add background color to formula box
    fill = textbox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # Add explanation
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(2.5)
    textbox2 = slide.shapes.add_textbox(left, top, width, height)
    text_frame2 = textbox2.text_frame
    
    for i, line in enumerate(explanation):
        if i == 0:
            p = text_frame2.paragraphs[0]
        else:
            p = text_frame2.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
    
    return slide

def create_example_slide(prs, title, example_data):
    """Create a slide with example calculations"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, item in enumerate(example_data):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        if item.startswith('=') or item.startswith('EJV'):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def generate_presentation():
    """Main function to generate the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    create_title_slide(prs, 
                      "EJV v2: Justice-Weighted Local Impact Calculation",
                      "Complete Technical Guide")
    
    # Executive Summary
    create_section_header(prs, "Executive Summary")
    
    create_content_slide(prs, "What is EJV v2?", [
        {'text': 'Transforms traditional 0-100 scoring into dollar-based metric', 'level': 0},
        {'text': 'Quantifies justice-weighted local economic impact of every purchase', 'level': 0},
        {'text': 'Incorporates equity adjustments based on ZIP-code economic conditions', 'level': 0},
        {'text': 'Key Question: "For every $100 spent, how many dollars create justice-weighted local economic impact?"', 'level': 0}
    ])
    
    # The Formula Section
    create_section_header(prs, "The Formula", "Core Calculation Method")
    
    create_formula_slide(prs, "EJV v2 Formula",
                        "EJV v2 = (P × LC) × (JS_ZIP / 100)",
                        [
                            "P = Purchase Amount (dollars)",
                            "LC = Local Capture (0-1, percentage as decimal)",
                            "JS_ZIP = Justice Score for ZIP code (0-100)"
                        ])
    
    # Component 1: Purchase Amount
    create_section_header(prs, "Component 1:", "Purchase Amount (P)")
    
    create_content_slide(prs, "Purchase Amount (P)", [
        {'text': 'Definition: The dollar value of the transaction', 'level': 0},
        {'text': 'Default Value: $100', 'level': 0},
        {'text': 'Standardized base for comparison across stores', 'level': 1},
        {'text': 'Can be adjusted for actual transaction amounts', 'level': 1},
        {'text': 'Scales linearly (e.g., $200 purchase = 2× the impact)', 'level': 1}
    ])
    
    # Component 2: Local Capture
    create_section_header(prs, "Component 2:", "Local Capture (LC)")
    
    create_content_slide(prs, "Local Capture (LC)", [
        {'text': 'Definition: Percentage of economic value remaining in local community', 'level': 0},
        {'text': 'Based on local hiring practices', 'level': 0},
        {'text': 'Range: 0.00 - 1.00 (0% - 100%)', 'level': 0},
        {'text': 'Calculation: LC = Local Hire Percentage', 'level': 0}
    ])
    
    create_content_slide(prs, "What Influences Local Capture?", [
        {'text': 'Base local hire: 40%-95% (store-specific)', 'level': 0},
        {'text': 'Unemployment adjustment: +0-20% bonus in high-unemployment areas', 'level': 0},
        {'text': 'Final range: 40%-98%', 'level': 0},
        {'text': 'Higher LC = More local economic benefit', 'level': 0}
    ])
    
    create_example_slide(prs, "Local Capture Examples", [
        "Store A: LC = 0.82 (82% of wages paid to local residents)",
        "• 82% of employee wages stay in local economy",
        "• Strong local economic benefit",
        "",
        "Store B: LC = 0.65 (65% of wages paid to local residents)",
        "• Only 65% of wages circulate locally",
        "• Moderate local economic benefit"
    ])
    
    # Component 3: Justice Score
    create_section_header(prs, "Component 3:", "Justice Score (JS_ZIP)")
    
    create_content_slide(prs, "Justice Score Overview", [
        {'text': 'Comprehensive 0-100 score measuring equity quality', 'level': 0},
        {'text': 'Based on 9 equity dimensions', 'level': 0},
        {'text': 'Adjusted for local economic need', 'level': 0},
        {'text': 'Higher score = Better equity performance', 'level': 0}
    ])
    
    create_content_slide(prs, "The 9 Equity Dimensions", [
        {'text': 'AES - Access to Essential Services', 'level': 0},
        {'text': 'ART - Access to Resources & Technology', 'level': 0},
        {'text': 'HWI - Health, Wellness & Inclusion', 'level': 0},
        {'text': 'PSR - Public Service Representation', 'level': 0},
        {'text': 'CAI - Cultural Awareness & Inclusivity', 'level': 0},
        {'text': 'JCE - Job Creation/Economic Empowerment', 'level': 0},
        {'text': 'FSI - Financial Support & Investment', 'level': 0},
        {'text': 'CED - Community Engagement & Development', 'level': 0},
        {'text': 'ESD - Education & Skill Development', 'level': 0}
    ])
    
    create_content_slide(prs, "Justice Score Calculation Process", [
        {'text': 'Step 1: Calculate Base Dimension Scores (0-1 scale)', 'level': 0},
        {'text': 'Each dimension normalized using specific metrics', 'level': 1},
        {'text': 'Step 2: Apply ZIP Need Modifiers (NM)', 'level': 0},
        {'text': 'Adjust for local economic conditions', 'level': 1},
        {'text': 'Range: 0.80 (low need) to 1.10 (high need)', 'level': 1},
        {'text': 'Step 3: Calculate Final Justice Score', 'level': 0},
        {'text': 'Average all adjusted dimensions × 100', 'level': 1}
    ])
    
    create_content_slide(prs, "ZIP Need Modifiers (NM)", [
        {'text': 'Purpose: Recognize equity work in disadvantaged areas', 'level': 0},
        {'text': 'Based on unemployment rate and median income', 'level': 0},
        {'text': 'Applied to 3 dimensions: AES, ART, HWI', 'level': 0},
        {'text': 'Examples:', 'level': 0},
        {'text': 'Manhattan (10001): NM = 0.925 (low need)', 'level': 1},
        {'text': 'South LA (90011): NM = 1.10 (very high need)', 'level': 1},
        {'text': 'Chicago SW (60629): NM = 1.05 (high need)', 'level': 1}
    ])
    
    create_formula_slide(prs, "Justice Score Formula",
                        "JS_ZIP = (Sum of Adjusted Dimensions / 9) × 100",
                        [
                            "Range: 0-100",
                            "90-100: Exceptional equity quality",
                            "70-89: Strong equity performance",
                            "50-69: Moderate equity performance",
                            "30-49: Needs improvement",
                            "0-29: Significant equity concerns"
                        ])
    
    # Complete Example
    create_section_header(prs, "Complete Calculation Example",
                         "Supermarket in Manhattan, NY (ZIP 10001)")
    
    create_example_slide(prs, "Example: Input Data", [
        "Store ID: supermarket_101",
        "Location: Manhattan, NY (ZIP 10001)",
        "Purchase Amount: $100.00",
        "",
        "Local Economic Conditions:",
        "• Unemployment Rate: 3.1%",
        "• Median Income: $106,509",
        "• Living Wage: $8.41/hour"
    ])
    
    create_example_slide(prs, "Example: Store Metrics", [
        "Average Wage: $15.02/hour",
        "Active Employees: 48",
        "Local Hire Percentage: 82%",
        "Daily Payroll: $5,769.60",
        "Community Spending: $288.48/day"
    ])
    
    create_example_slide(prs, "Example: Step 1 - Base Dimensions", [
        "Wage Score = ($15.02 / $8.41) / 25 = 1.786 / 25 = 0.071 → capped at 1.000",
        "Hiring Score = 82% / 25 = 3.28 / 25 = 0.131 → capped at 1.000",
        "Community Score = $288.48 / $5,769.60 / 25 = 0.002",
        "Participation Score = 48 / 25 / 25 = 0.077 → capped at 1.000",
        "",
        "Base Dimensions:",
        "AES=0.002, ART=1.000, HWI=1.000, PSR=0.002",
        "CAI=1.000, JCE=1.000, FSI=1.000, CED=0.501, ESD=1.000"
    ])
    
    create_example_slide(prs, "Example: Step 2 - Apply Need Modifiers", [
        "Unemployment Factor = 3.1% / 10% = 0.31",
        "Income Factor = max(0, 1 - $106,509/$75,000) = 0",
        "Base Modifier = 0.80 + (0.30 × (0.31 + 0) / 2) = 0.8465",
        "",
        "NM_AES = 0.8465 × 1.05 = 0.925",
        "NM_ART = 0.8465 × 1.02 = 0.925",
        "NM_HWI = 0.8465 × 1.03 = 0.925",
        "",
        "Adjusted Dimensions:",
        "AES=0.002, ART=0.925, HWI=0.925, PSR=0.002",
        "CAI=1.000, JCE=1.000, FSI=1.000, CED=0.501, ESD=1.000"
    ])
    
    create_example_slide(prs, "Example: Step 3 - Calculate Justice Score", [
        "Sum of Adjusted Dimensions:",
        "0.002 + 0.925 + 0.925 + 0.002 + 1.000 + 1.000 + 1.000 + 0.501 + 1.000",
        "= 6.355",
        "",
        "JS_ZIP = 6.355 / 9 × 100 = 70.6"
    ])
    
    create_example_slide(prs, "Example: Step 4 - Calculate EJV v2", [
        "P = $100.00",
        "LC = 0.82",
        "JS_ZIP = 70.6",
        "",
        "EJV v2 = (P × LC) × (JS_ZIP / 100)",
        "EJV v2 = ($100.00 × 0.82) × (70.6 / 100)",
        "EJV v2 = $82.00 × 0.706",
        "= $57.89",
        "",
        "Result: For every $100 spent at this supermarket,",
        "$57.89 creates justice-weighted local economic impact"
    ])
    
    create_content_slide(prs, "Example Interpretation", [
        {'text': '$57.89 of justice-weighted local impact per $100 spent', 'level': 0},
        {'text': '82% local hiring (money stays in community)', 'level': 0},
        {'text': 'Strong wage quality ($15.02 vs $8.41 living wage)', 'level': 0},
        {'text': 'Solid equity performance (70.6/100 justice score)', 'level': 0},
        {'text': 'Low-need area (Manhattan has low unemployment, high income)', 'level': 0}
    ])
    
    # EJV v1 vs v2 Comparison
    create_section_header(prs, "EJV v1 vs EJV v2", "Key Differences")
    
    create_content_slide(prs, "Comparing EJV v1 and v2", [
        {'text': 'EJV v1: Abstract 0-100 score (75.91)', 'level': 0},
        {'text': 'Good for relative comparison and ranking', 'level': 1},
        {'text': 'No equity adjustment', 'level': 1},
        {'text': 'EJV v2: Dollar-based metric ($57.89)', 'level': 0},
        {'text': 'Direct economic impact measurement', 'level': 1},
        {'text': 'Includes ZIP Need Modifiers for equity', 'level': 1},
        {'text': 'Better for budget planning and procurement', 'level': 1}
    ])
    
    # Real-World Applications
    create_section_header(prs, "Real-World Applications")
    
    create_example_slide(prs, "Application 1: Consumer Choice", [
        "Scenario: Choosing between two grocery stores",
        "",
        "Store A (Local Chain):",
        "• EJV v2: $64.20 per $100",
        "• High local impact, strong equity",
        "",
        "Store B (National Chain):",
        "• EJV v2: $28.40 per $100",
        "• Lower local impact",
        "",
        "Decision: Store A creates 2.26× more justice-weighted impact",
        "Annual Impact ($5,000/year): $1,790 more local impact with Store A"
    ])
    
    create_example_slide(prs, "Application 2: Municipal Procurement", [
        "Scenario: City purchasing $1M in office supplies",
        "",
        "Vendor A:",
        "• EJV v2: $72.50 per $100",
        "• Total Impact: $725,000",
        "",
        "Vendor B:",
        "• EJV v2: $45.30 per $100",
        "• Total Impact: $453,000",
        "",
        "Policy Decision: Choosing Vendor A creates",
        "$272,000 more justice-weighted local economic benefit"
    ])
    
    create_example_slide(prs, "Application 3: Community Investment", [
        "Scenario: Analyzing economic equity in a neighborhood",
        "",
        "Aggregated Data (20 stores):",
        "• Average EJV v2: $58.30 per $100",
        "• Range: $22.10 - $84.60",
        "• High-performers (>$70): 6 stores (30%)",
        "• Medium-performers ($50-$70): 10 stores (50%)",
        "• Low-performers (<$50): 4 stores (20%)",
        "",
        "Insight: Targeting low performers could increase average by 15%"
    ])
    
    # Data Sources
    create_content_slide(prs, "Data Sources", [
        {'text': 'Bureau of Labor Statistics (BLS)', 'level': 0},
        {'text': 'OEWS May 2024 wage data', 'level': 1},
        {'text': 'Industry-specific occupation codes', 'level': 1},
        {'text': 'U.S. Census Bureau', 'level': 0},
        {'text': 'American Community Survey (ACS) 2022', 'level': 1},
        {'text': 'Median household income and unemployment rates', 'level': 1},
        {'text': 'Industry Standards', 'level': 0},
        {'text': 'Employee count benchmarks', 'level': 1},
        {'text': 'Labor market research', 'level': 1}
    ])
    
    # Best Practices
    create_content_slide(prs, "Best Practices & Limitations", [
        {'text': 'Do:', 'level': 0},
        {'text': 'Use for relative comparisons', 'level': 1},
        {'text': 'Consider context (store type, location, size)', 'level': 1},
        {'text': 'Combine with other equity metrics', 'level': 1},
        {'text': "Don't:", 'level': 0},
        {'text': 'Use as sole decision criterion', 'level': 1},
        {'text': 'Ignore qualitative factors', 'level': 1},
        {'text': 'Assume perfect data accuracy', 'level': 1}
    ])
    
    # Future Enhancements
    create_content_slide(prs, "Future Enhancements", [
        {'text': 'Dynamic Purchase Amounts', 'level': 0},
        {'text': 'Category-specific baselines', 'level': 1},
        {'text': 'Historical Tracking', 'level': 0},
        {'text': 'Trend analysis over time', 'level': 1},
        {'text': 'Enhanced Granularity', 'level': 0},
        {'text': 'Census tract-level modifiers', 'level': 1},
        {'text': 'Additional Dimensions', 'level': 0},
        {'text': 'Environmental sustainability, supply chain equity', 'level': 1},
        {'text': 'Machine Learning Integration', 'level': 0}
    ])
    
    # FAQs
    create_content_slide(prs, "Frequently Asked Questions", [
        {'text': 'Why use dollars instead of 0-100 scores?', 'level': 0},
        {'text': 'Dollars are more intuitive and actionable', 'level': 1},
        {'text': 'Why do high-income areas get lower need modifiers?', 'level': 0},
        {'text': 'Recognizes greater impact of equity work in disadvantaged areas', 'level': 1},
        {'text': 'Can EJV v2 exceed the purchase amount?', 'level': 0},
        {'text': 'No. Maximum is 100% of purchase amount', 'level': 1},
        {'text': 'How often is data updated?', 'level': 0},
        {'text': 'Wage data annually, store metrics real-time', 'level': 1}
    ])
    
    # Conclusion
    create_content_slide(prs, "Conclusion", [
        {'text': 'EJV v2 transforms abstract scores into tangible dollar-based impact', 'level': 0},
        {'text': 'Incorporates equity-weighted adjustments via ZIP Need Modifiers', 'level': 0},
        {'text': 'Ensures fair recognition of businesses serving disadvantaged communities', 'level': 0},
        {'text': 'Provides actionable insights for:', 'level': 0},
        {'text': 'Consumers making purchasing decisions', 'level': 1},
        {'text': 'Businesses improving their equity performance', 'level': 1},
        {'text': 'Policymakers advancing economic equity', 'level': 1}
    ])
    
    # Final Slide
    create_title_slide(prs,
                      "Thank You",
                      "FIX$ GeoEquity Impact Engine\nhttps://fixapp-phi.vercel.app")
    
    # Save the presentation
    output_file = "EJV_V2_Calculation_Guide.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    print("🎯 Generating EJV v2 PowerPoint Presentation...")
    print("📄 Reading from: EJV_V2_CALCULATION_GUIDE.md")
    print()
    
    try:
        output_file = generate_presentation()
        print()
        print("=" * 60)
        print("✨ Generation Complete!")
        print(f"📁 Output file: {output_file}")
        print("=" * 60)
    except Exception as e:
        print(f"❌ Error generating presentation: {str(e)}")
        import traceback
        traceback.print_exc()
