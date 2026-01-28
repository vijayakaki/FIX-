"""
FIX$ Application - PowerPoint Presentation Generator
Generates a complete .pptx PowerPoint presentation with file descriptions and comparisons
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    return slide

def create_content_slide(prs, title, content_items):
    """Create a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item['file']
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add description line 1
        p1 = text_frame.add_paragraph()
        p1.text = item['desc_line1']
        p1.level = 1
        p1.font.size = Pt(14)
        p1.font.color.rgb = RGBColor(60, 60, 60)
        
        # Add description line 2
        p2 = text_frame.add_paragraph()
        p2.text = item['desc_line2']
        p2.level = 1
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(60, 60, 60)
    
    return slide

def create_comparison_slide(prs, comparison_title, comparison_items):
    """Create a comparison slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = comparison_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 126, 234)
    p.alignment = PP_ALIGN.CENTER
    
    # Add comparison content
    content_top = Inches(1.5)
    content_left = Inches(0.5)
    content_width = Inches(9)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for key, value in comparison_items.items():
        # Add key (bold)
        p = text_frame.add_paragraph()
        p.text = f"{key}:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 126, 234)
        p.space_after = Pt(6)
        
        # Add value
        p = text_frame.add_paragraph()
        p.text = value
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(18)
        p.level = 1
    
    return slide

def generate_presentation():
    """Generate complete PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    create_title_slide(
        prs,
        "FIX$ Application File Structure",
        "Complete Guide to All Application Files"
    )
    
    # File descriptions by category
    file_descriptions = {
        "Core Backend Files": [
            {
                "file": "app.py",
                "desc_line1": "Flask application server for local development with all API endpoints and calculation logic.",
                "desc_line2": "Runs on localhost:5000 when you execute 'python run.py' for testing on your computer."
            },
            {
                "file": "api/index.py",
                "desc_line1": "Production version of Flask app deployed on Vercel cloud platform as serverless function.",
                "desc_line2": "Same code as app.py but optimized for cloud deployment at https://fix-app-three.vercel.app"
            },
            {
                "file": "database.py",
                "desc_line1": "SQLite database manager handling user accounts, authentication, and session management.",
                "desc_line2": "Uses file-based storage locally (fixapp.db) and in-memory storage on Vercel."
            },
            {
                "file": "run.py",
                "desc_line1": "Development server launcher that initializes database and starts Flask app on localhost:5000.",
                "desc_line2": "Provides console messages with URLs, demo credentials, and enables debug mode."
            }
        ],
        "Frontend Files": [
            {
                "file": "public/index.html",
                "desc_line1": "Main application interface with interactive map (Leaflet.js), store search, and EJV visualization.",
                "desc_line2": "Communicates with backend APIs to display economic justice calculations."
            },
            {
                "file": "public/login-simple.html",
                "desc_line1": "User authentication page with username/password input and session management.",
                "desc_line2": "Redirects to main app after successful login (demo: admin/fix123)."
            }
        ],
        "Configuration Files": [
            {
                "file": "vercel.json",
                "desc_line1": "Deployment configuration for Vercel cloud platform specifying routes and build settings.",
                "desc_line2": "Routes API requests to index.py (Python serverless) and static files to public/."
            },
            {
                "file": "requirements.txt",
                "desc_line1": "Python package dependencies list: Flask, Flask-CORS, requests, Werkzeug, python-pptx.",
                "desc_line2": "Install all dependencies with 'pip install -r requirements.txt' command."
            },
            {
                "file": ".gitignore",
                "desc_line1": "Specifies files/folders to exclude from Git version control (databases, cache, secrets).",
                "desc_line2": "Prevents sensitive data (.env files) and build artifacts from being committed."
            }
        ],
        "Documentation Files": [
            {
                "file": "EJV_COMPLETE_CALCULATION_GUIDE.md",
                "desc_line1": "Comprehensive guide with all 6 EJV calculation versions, formulas, and worked examples.",
                "desc_line2": "Includes data sources (BLS, Census, BEA), API documentation, and citations."
            },
            {
                "file": "FILE_STRUCTURE_DOCUMENTATION.md",
                "desc_line1": "Complete guide explaining every file in the application with purpose and features.",
                "desc_line2": "Technical reference for developers understanding the codebase structure."
            },
            {
                "file": "LINKEDIN_POST.md",
                "desc_line1": "Social media content showcasing AI-powered development journey.",
                "desc_line2": "Highlights CMMC AI, Medical AI, Urban AI, and FIX$ apps with tech stack."
            }
        ],
        "Database & Cache Files": [
            {
                "file": "fixapp.db",
                "desc_line1": "SQLite database file storing user accounts and sessions (local development only).",
                "desc_line2": "Auto-generated when running locally, not used on Vercel (uses in-memory)."
            },
            {
                "file": "__pycache__/",
                "desc_line1": "Python bytecode cache folder for faster loading of compiled .pyc files.",
                "desc_line2": "Auto-generated by Python, safe to delete (regenerates automatically)."
            }
        ]
    }
    
    # Create slides for each category
    for category, files in file_descriptions.items():
        create_content_slide(prs, category, files)
    
    # Comparison slides
    comparisons = {
        "app.py vs api/index.py": {
            "app.py": "Local development - test on your computer (localhost:5000)",
            "api/index.py": "Production deployment - live on internet via Vercel cloud"
        },
        "Flask vs Vercel": {
            "Flask": "Python library that handles requests at API endpoints, runs functions, returns JSON",
            "Vercel": "Cloud hosting platform that deploys your code to internet as serverless application"
        },
        "Vercel vs Bluehost": {
            "Vercel": "Serverless (pay-per-use, auto-scaling, GitHub integration)",
            "Bluehost": "Traditional hosting (always-on server, fixed resources, manual FTP)"
        },
        "Frontend vs Backend": {
            "Frontend (HTML/JS)": "What users see - interactive map, forms, buttons, visualization",
            "Backend (Python/Flask)": "What users don't see - calculations, database, API logic"
        }
    }
    
    # Add comparison slides
    for comparison_title, comparison_items in comparisons.items():
        create_comparison_slide(prs, comparison_title, comparison_items)
    
    # Final slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 126, 234)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = "FIX$ GeoEquity Impact Engine"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    filename = 'FIX_APP_FILE_STRUCTURE.pptx'
    prs.save(filename)
    print(f"\n✅ PowerPoint presentation created: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"💡 Open with Microsoft PowerPoint or Google Slides\n")
    
    return filename

if __name__ == "__main__":
    print("\n" + "="*70)
    print("FIX$ APPLICATION - POWERPOINT PRESENTATION GENERATOR")
    print("="*70 + "\n")
    
    print("🎯 Generating PowerPoint presentation...\n")
    
    try:
        filename = generate_presentation()
        print("="*70)
        print("✅ SUCCESS! Your presentation is ready.")
        print("="*70)
        print(f"\n📁 File Location: {filename}")
        print("\n📌 What's included:")
        print("   • Title slide")
        print("   • Core Backend Files (4 files)")
        print("   • Frontend Files (2 files)")
        print("   • Configuration Files (3 files)")
        print("   • Documentation Files (3 files)")
        print("   • Database & Cache Files (2 files)")
        print("   • 4 Comparison slides")
        print("   • Thank you slide")
        print("\n💡 Tips:")
        print("   • Edit styles and colors as needed")
        print("   • Add images or diagrams")
        print("   • Customize for your audience")
        print("\n")
        
    except ImportError:
        print("❌ ERROR: python-pptx library not installed")
        print("\n📦 Install with:")
        print("   pip install python-pptx")
        print("\nThen run this script again.")
    except Exception as e:
        print(f"❌ ERROR: {e}")
        print("\nPlease check the error message and try again.")
